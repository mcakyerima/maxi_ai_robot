# -*- coding: utf-8 -*-
"""
Maxi - AI Educational Robot
Presentation builder for the bootcamp showcase (4 August 2026).

Run:  python build_presentation.py
Out:  Maxi_Robot_Presentation.pptx

Design system is shared VERBATIM with the sibling deck (Brain Tumour Detection)
so the two look like siblings when presented back to back. Only the content is
Maxi's. Edit the CONFIG block below to set names and the fill-in hardware BOM.

Photos next to this file: robot_1.png (hero), robot_2.png, robot_3.png — each
add_picture is guarded with os.path.exists() so a missing file never crashes.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# --------------------------------------------------------------------------
# CONFIG  <-- edit these
# --------------------------------------------------------------------------
PRODUCT    = "Maxi"
BOOTCAMP   = "6-Month Coding Bootcamp"
COHORT     = "Secondary School Cohort"
EVENT_DATE = "4 August 2026"
PRESENTERS = ["Student One", "Student Two", "Student Three", "Student Four"]

HERO    = "robot_1.png"   # title hero shot
PHOTO_2 = "robot_2.png"   # the build slide
PHOTO_3 = "robot_3.png"   # what-it-can-do / demo

# Bill of materials — NOT in the code. Fill these in with the real values before
# printing; do not guess on stage. Leaving them shows a clear "confirm" placeholder.
HARDWARE = {
    "pi":     "Raspberry Pi  —  (add exact model)",
    "servos": "12 × micro-servos  —  (add model, e.g. SG90 / MG90S)",
    "power":  "Servo power  —  (add source, e.g. external 5V supply)",
    "cost":   "(add total build cost)",
    "parts":  "(add part count)",
    "body":   "(3D-printed? add material)",
}

# --------------------------------------------------------------------------
# Design system  (shared verbatim with the sibling deck)
# --------------------------------------------------------------------------
W, H = Inches(13.333), Inches(7.5)

BG_A     = RGBColor(0x06, 0x09, 0x11)   # deep space
BG_B     = RGBColor(0x0D, 0x14, 0x22)   # slate navy
PANEL    = RGBColor(0x12, 0x1A, 0x28)
PANEL_2  = RGBColor(0x18, 0x22, 0x33)
STROKE   = RGBColor(0x27, 0x37, 0x50)

PURPLE   = RGBColor(0x8B, 0x7F, 0xF5)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
GREEN    = RGBColor(0x2D, 0xF5, 0x9B)
PINK     = RGBColor(0xF4, 0x72, 0xB6)
AMBER    = RGBColor(0xFB, 0xBF, 0x24)
RED      = RGBColor(0xF8, 0x71, 0x71)

INK      = RGBColor(0xF3, 0xF7, 0xFF)
MUTED    = RGBColor(0x93, 0xA5, 0xC0)
DIM      = RGBColor(0x64, 0x77, 0x94)

FONT   = "Segoe UI"
FONT_B = "Segoe UI Semibold"
MONO   = "Consolas"

ACCENTS = [PURPLE, CYAN, GREEN, PINK, AMBER]


# --------------------------------------------------------------------------
# Low-level helpers  (verbatim)
# --------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()
    return shape


def _no_shadow(shape):
    shape.shadow.inherit = False
    return shape


def set_alpha(shape, pct):
    """Apply transparency (0-100) to a solid-filled shape."""
    xPr = shape.fill._xPr
    solid = xPr.find(qn("a:solidFill"))
    if solid is None:
        return shape
    clr = solid.find(qn("a:srgbClr"))
    if clr is None:
        return shape
    a = parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="%d"/>'
                  % int((100 - pct) * 1000))
    clr.append(a)
    return shape


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, alpha=0, radius=None):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if alpha:
        set_alpha(s, alpha)
    _no_line(s)
    _no_shadow(s)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    s.text_frame.word_wrap = True
    return s


def panel(slide, x, y, w, h, fill=PANEL, border=STROKE, radius=0.055, alpha=0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if alpha:
        set_alpha(s, alpha)
    s.line.color.rgb = border
    s.line.width = Pt(1)
    _no_shadow(s)
    try:
        s.adjustments[0] = radius
    except (IndexError, ValueError):
        pass
    s.text_frame.word_wrap = True
    return s


def text(slide, x, y, w, h, content, size=18, color=INK, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, spacing=1.0, space_after=0, anchor=MSO_ANCHOR.TOP,
         italic=False):
    """content: str with \n for new paragraphs, or list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = content.split("\n") if isinstance(content, str) else content
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        if isinstance(line, tuple):
            txt, sz, col, bd = (list(line) + [size, color, bold])[:4]
        else:
            txt, sz, col, bd = line, size, color, bold
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = bd
        r.font.italic = italic
        r.font.name = FONT_B if bd and font == FONT else font
    return tb


def glow(slide, cx, cy, r, color, alpha=88):
    """Soft neon orb for atmosphere."""
    s = rect(slide, cx - r, cy - r, 2 * r, 2 * r, color, MSO_SHAPE.OVAL, alpha=alpha)
    return s


def background(slide, orbs=((0.16, 0.14, PURPLE), (0.88, 0.82, CYAN))):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.gradient()
    stops = bg.fill.gradient_stops
    stops[0].color.rgb = BG_A
    stops[0].position = 0.0
    stops[1].color.rgb = BG_B
    stops[1].position = 1.0
    bg.fill.gradient_angle = 315.0
    _no_line(bg)
    _no_shadow(bg)
    for fx, fy, col in orbs:
        glow(slide, Emu(int(W * fx)), Emu(int(H * fy)), Inches(2.4), col, alpha=90)
    return bg


def badge(slide, x, y, size, emoji, color, glyph_size=None):
    """Circular icon chip with an emoji glyph."""
    ring = rect(slide, x, y, size, size, color, MSO_SHAPE.OVAL, alpha=80)
    inner_pad = Emu(int(size * 0.12))
    core = rect(slide, x + inner_pad, y + inner_pad,
                size - 2 * inner_pad, size - 2 * inner_pad, color, MSO_SHAPE.OVAL, alpha=62)
    tf = core.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = emoji
    r.font.size = Pt(glyph_size or int(size / Inches(1) * 22))
    r.font.name = "Segoe UI Emoji"
    r.font.color.rgb = INK
    return ring


def pill(slide, x, y, label, color, w=None, size=11):
    w = w or Inches(0.16 * len(label) + 0.34)
    s = rect(slide, x, y, w, Inches(0.3), color, MSO_SHAPE.ROUNDED_RECTANGLE, alpha=76, radius=0.5)
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT_B
    r.font.color.rgb = color
    return s


def notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body


def framed_picture(slide, path, x, y, w, max_h, border=PURPLE, placeholder="[ photo ]"):
    """Rounded panel behind a photo, aspect preserved, scaled to fit (w × max_h),
    horizontally centred. Guarded so a missing file just draws a labelled panel."""
    if os.path.exists(path):
        outer = panel(slide, x, y, w, max_h + Inches(0.24), fill=BG_A, border=border)
        pic = slide.shapes.add_picture(path, x + Inches(0.12), y + Inches(0.12),
                                       width=w - Inches(0.24))
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.height = int(max_h)
            pic.width = int(pic.width * ratio)
            pic.left = int(x + (w - pic.width) / 2)
        outer.height = int(pic.height + Inches(0.24))
        return outer
    ph = panel(slide, x, y, w, max_h, fill=BG_A, border=border)
    text(slide, x, y + Emu(int(max_h / 2)) - Inches(0.3), w, Inches(0.6), placeholder,
         size=13, color=DIM, align=PP_ALIGN.CENTER, spacing=1.3)
    return ph


# --------------------------------------------------------------------------
# Slide scaffolding
# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

_slide_no = {"n": 0}


def new_slide(orbs=((0.16, 0.14, PURPLE), (0.88, 0.82, CYAN))):
    s = prs.slides.add_slide(BLANK)
    background(s, orbs)
    return s


def header(slide, kicker, title, subtitle=None, accent=PURPLE, number=None):
    _slide_no["n"] += 1
    n = number if number is not None else _slide_no["n"]

    rail = rect(slide, Inches(0), Inches(0), Inches(0.075), H, accent)
    set_alpha(rail, 25)

    pill(slide, Inches(0.62), Inches(0.46), kicker.upper(), accent)
    text(slide, Inches(0.62), Inches(0.86), Inches(11.4), Inches(0.75), title,
         size=34, bold=True, color=INK)
    y = Inches(1.62)
    if subtitle:
        text(slide, Inches(0.62), Inches(1.55), Inches(11.4), Inches(0.4), subtitle,
             size=14, color=MUTED)
        y = Inches(2.02)
    underline = rect(slide, Inches(0.62), y - Inches(0.12), Inches(1.5), Inches(0.045), accent)
    set_alpha(underline, 10)

    text(slide, Inches(0.62), Inches(6.95), Inches(9.0), Inches(0.3),
         "%s  ·  %s  ·  %s" % (PRODUCT, BOOTCAMP, EVENT_DATE), size=9, color=DIM)
    text(slide, Inches(11.6), Inches(6.95), Inches(1.1), Inches(0.3), "%02d" % n,
         size=9, color=DIM, align=PP_ALIGN.RIGHT)
    return y + Inches(0.24)


def footer_only(slide):
    text(slide, Inches(0.62), Inches(6.95), Inches(9.0), Inches(0.3),
         "%s  ·  %s  ·  %s" % (PRODUCT, BOOTCAMP, EVENT_DATE), size=9, color=DIM)
    text(slide, Inches(11.6), Inches(6.95), Inches(1.1), Inches(0.3),
         "%02d" % _slide_no["n"], size=9, color=DIM, align=PP_ALIGN.RIGHT)


# ==========================================================================
# 01 - TITLE
# ==========================================================================
s = new_slide(orbs=((0.2, 0.2, PURPLE), (0.82, 0.78, CYAN), (0.55, 0.9, PINK)))

# left: identity
text(s, Inches(0.85), Inches(1.55), Inches(6.3), Inches(0.4),
     "MEET  MAXI", size=15, bold=True, color=CYAN)
text(s, Inches(0.8), Inches(2.0), Inches(6.4), Inches(1.2), "Maxi",
     size=76, bold=True, color=INK)
text(s, Inches(0.85), Inches(3.35), Inches(6.2), Inches(0.5),
     "A robot that listens, talks back,\nand teaches a child on its fingers",
     size=19, color=CYAN, spacing=1.2)
text(s, Inches(0.85), Inches(4.35), Inches(6.1), Inches(0.4),
     "Tablet face  ·  cloud brain  ·  two robotic hands", size=13, color=MUTED)

chip_x = Inches(0.85)
for lbl, col, wd in [("Python", CYAN, 1.0), ("Flask + Socket.IO", PURPLE, 1.75),
                     ("Groq LLM", GREEN, 1.15), ("Raspberry Pi", PINK, 1.45)]:
    pill(s, chip_x, Inches(4.95), lbl, col, w=Inches(wd), size=11)
    chip_x += Inches(wd) + Inches(0.16)

text(s, Inches(0.85), Inches(5.7), Inches(6.2), Inches(0.9),
     "%s  ·  %s\nPresented by %s\n%s"
     % (BOOTCAMP, COHORT, ", ".join(PRESENTERS), EVENT_DATE),
     size=12, color=MUTED, spacing=1.35)

# right: hero photo, framed and fit
framed_picture(s, HERO, Inches(7.45), Inches(1.55), Inches(5.25), Inches(4.35),
               border=PURPLE, placeholder="[ robot_1.png — hero shot lands here ]")

_slide_no["n"] = 1
notes(s, """[SPEAKER 1 — 10 sec]
"Good morning. We are the %s of the %s. Over six months we built Maxi — a small
robot that listens to a child, talks back in a friendly voice, and even counts
maths out on real robotic fingers."
Gesture to the robot on the table. Hold a beat, then advance.""" % (COHORT, BOOTCAMP))


# ==========================================================================
# 02 - AGENDA / TEAM
# ==========================================================================
s = new_slide()
top = header(s, "Agenda", "How the next 10 minutes go", accent=CYAN)

items = [
    ("😀", "Meet Maxi", "What it is and why we built it", PURPLE, PRESENTERS[0], "2 min"),
    ("🛠️", "The build & the brain", "Hands, board, and the cloud that thinks", CYAN,
     PRESENTERS[1] if len(PRESENTERS) > 1 else PRESENTERS[0], "2 min"),
    ("🧠", "How it thinks", "The state machine, the code, the wiring", GREEN,
     PRESENTERS[2] if len(PRESENTERS) > 2 else PRESENTERS[0], "2 min"),
    ("🚀", "Live demo", "We talk to Maxi and its fingers move", PINK,
     PRESENTERS[3] if len(PRESENTERS) > 3 else PRESENTERS[0], "3 min"),
    ("💡", "Lessons & limits", "What broke, what we learned, what's next", AMBER,
     PRESENTERS[0], "1 min"),
]
y = top
for emo, title_, desc, col, who, mins in items:
    panel(s, Inches(0.62), y, Inches(12.1), Inches(0.78), fill=PANEL, alpha=18)
    badge(s, Inches(0.85), y + Inches(0.15), Inches(0.48), emo, col, glyph_size=15)
    text(s, Inches(1.5), y + Inches(0.13), Inches(3.2), Inches(0.3), title_, size=15, bold=True)
    text(s, Inches(1.5), y + Inches(0.43), Inches(6.6), Inches(0.3), desc, size=11.5, color=MUTED)
    text(s, Inches(8.35), y + Inches(0.25), Inches(2.75), Inches(0.3), who, size=11.5, color=col,
         align=PP_ALIGN.RIGHT, bold=True)
    pill(s, Inches(11.45), y + Inches(0.24), mins, col, w=Inches(0.95), size=10)
    y += Inches(0.9)

text(s, Inches(0.62), Inches(6.5), Inches(12.1), Inches(0.35),
     "We each built a different part, so we will each explain our own part.",
     size=12.5, color=MUTED, italic=True)

notes(s, """[SPEAKER 1 — 25 sec]
Point at each row as you name the teammate who owns it, so the audience knows who
speaks when. Keep it to one line per row — do not read the descriptions out.
"We split into hands, brain and app, so you will hear from all four of us."
Then continue straight into the problem.""")


# ==========================================================================
# 03 - THE PROBLEM
# ==========================================================================
s = new_slide(orbs=((0.15, 0.2, RED), (0.85, 0.8, PURPLE)))
top = header(s, "01 · The problem", "One teacher, a whole classroom of children", accent=RED)

cards = [
    ("👩‍🏫", "Not enough one-on-one", "In a big class, a child who is stuck rarely gets a "
     "patient, private explanation. The quiet ones fall behind first.", PURPLE),
    ("✋", "Learning is physical", "A six-year-old understands \"three\" by seeing three "
     "fingers, not by reading it. Numbers need to be shown, not just said.", CYAN),
    ("🔁", "A robot never gets tired", "A child can ask the same question ten times. "
     "A patient machine answers the tenth time exactly like the first.", AMBER),
]
x = Inches(0.62)
for emo, t, d, col in cards:
    panel(s, x, top, Inches(3.86), Inches(3.05), alpha=14)
    badge(s, x + Inches(0.36), top + Inches(0.4), Inches(0.72), emo, col, glyph_size=21)
    text(s, x + Inches(0.36), top + Inches(1.4), Inches(3.2), Inches(0.4), t, size=17, bold=True)
    text(s, x + Inches(0.36), top + Inches(1.86), Inches(3.25), Inches(1.1), d,
         size=12, color=MUTED, spacing=1.3)
    x += Inches(4.04)

panel(s, Inches(0.62), Inches(5.42), Inches(12.1), Inches(1.15), fill=PANEL_2, border=PURPLE, alpha=12)
badge(s, Inches(0.95), Inches(5.7), Inches(0.58), "💡", PURPLE, glyph_size=17)
text(s, Inches(1.72), Inches(5.72), Inches(10.6), Inches(0.7),
     [("Our question:  ", 15, PURPLE, True),
      ("could we build a friendly robot that listens, talks back, and teaches a child "
       "patiently — one-on-one — using its own hands?", 15, INK, False)], spacing=1.2)

notes(s, """[SPEAKER 1 — 35 sec]
"Before the tech, here is why we built this."
One sentence per card — do NOT read them word for word. The middle card is the heart
of it: young kids learn numbers physically, on fingers. That is why Maxi has hands.
Land hard on the bottom line — that is the question the whole project answers.
IMPORTANT: we are not replacing teachers. Say it here; it earns trust for the rest.""")


# ==========================================================================
# 04 - MEET MAXI  (sense / think / act)
# ==========================================================================
s = new_slide()
top = header(s, "02 · Meet Maxi", "One loop: sense → think → act",
             "Everything Maxi does is this loop, running over and over", accent=GREEN)

loop = [
    ("👂", "SENSE", "It hears you", "The tablet microphone catches the child's voice and the "
     "wake word \"Hey Maxi\".", CYAN),
    ("🧠", "THINK", "It understands", "A cloud brain turns speech into an answer — chat, a maths "
     "step, a story — and decides what to do.", PURPLE),
    ("🗣️", "ACT", "It responds", "Maxi talks back in a warm voice, moves its 12 finger servos, "
     "and shows a feeling on its face.", GREEN),
]
x = Inches(0.62)
for i, (emo, tag, t, d, col) in enumerate(loop):
    panel(s, x, top, Inches(3.86), Inches(3.15), alpha=14)
    badge(s, x + Inches(0.36), top + Inches(0.38), Inches(0.8), emo, col, glyph_size=24)
    pill(s, x + Inches(1.42), top + Inches(0.62), tag, col, w=Inches(1.15), size=11)
    text(s, x + Inches(0.36), top + Inches(1.42), Inches(3.2), Inches(0.4), t, size=17, bold=True, color=col)
    text(s, x + Inches(0.36), top + Inches(1.9), Inches(3.25), Inches(1.15), d,
         size=12, color=MUTED, spacing=1.3)
    if i < 2:
        a = rect(s, x + Inches(3.9), top + Inches(1.35), Inches(0.24), Inches(0.34),
                 STROKE, MSO_SHAPE.CHEVRON)
        set_alpha(a, 45)
    x += Inches(4.04)

parts = [("📱", "Android tablet", "face, voice & ears", CYAN),
         ("🤲", "Two robotic hands", "12 servos, 10 fingers + 2 wrists", PURPLE),
         ("🍓", "Raspberry Pi", "drives the hands", PINK),
         ("☁️", "Cloud brain", "on Railway", GREEN)]
x = Inches(0.62)
for emo, t, d, col in parts:
    panel(s, x, Inches(5.42), Inches(2.9), Inches(1.15), fill=PANEL_2, alpha=16)
    badge(s, x + Inches(0.24), Inches(5.62), Inches(0.5), emo, col, glyph_size=15)
    text(s, x + Inches(0.86), Inches(5.6), Inches(1.95), Inches(0.32), t, size=12.5, bold=True, color=col)
    text(s, x + Inches(0.86), Inches(5.92), Inches(1.95), Inches(0.5), d, size=10, color=MUTED, spacing=1.1)
    x += Inches(3.06)

notes(s, """[SPEAKER 1 — 40 sec]
This is the mental model for the whole talk. Point across the three cards:
"Maxi senses with its ears, thinks in the cloud, and acts with its voice and hands.
Sense, think, act — that is the whole robot."
Then the bottom strip: the four physical pieces. "The face is a tablet. The hands are
real servos. A Raspberry Pi drives the hands. And the thinking happens in the cloud."
HAND OFF: "%s will show you how we actually built the body." """
      % (PRESENTERS[1] if len(PRESENTERS) > 1 else PRESENTERS[0]))


# ==========================================================================
# 05 - THE BUILD  (photo robot_2)
# ==========================================================================
s = new_slide(orbs=((0.18, 0.2, CYAN), (0.84, 0.78, PURPLE)))
top = header(s, "03 · The build", "What Maxi is made of", accent=CYAN)

# left: photo
framed_picture(s, PHOTO_2, Inches(0.62), top, Inches(5.55), Inches(3.35),
               border=CYAN, placeholder="[ robot_2.png — the build ]")

# right: real hardware facts from the code
facts = [
    ("🤲", "Two hands, 12 servos", "5 fingers + 1 wrist per hand. Servo channels 0–11 on one "
     "driver board.", PURPLE),
    ("🎛️", "Adafruit PCA9685", "A 16-channel PWM board that turns the Pi's commands into servo "
     "movement over I²C.", CYAN),
    ("📐", "Calibrated per finger", "Every finger has its own open/closed angles — e.g. right "
     "pinky 0–90°, wrist 20–160°.", GREEN),
]
y = top
for emo, t, d, col in facts:
    panel(s, Inches(6.42), y, Inches(6.3), Inches(1.05), fill=PANEL, alpha=16)
    badge(s, Inches(6.68), y + Inches(0.26), Inches(0.54), emo, col, glyph_size=15)
    text(s, Inches(7.4), y + Inches(0.15), Inches(5.1), Inches(0.32), t, size=14, bold=True, color=col)
    text(s, Inches(7.4), y + Inches(0.5), Inches(5.15), Inches(0.5), d, size=11, color=MUTED, spacing=1.2)
    y += Inches(1.18)

# bill of materials (fill-in)
panel(s, Inches(0.62), Inches(5.62), Inches(12.1), Inches(1.05), fill=PANEL_2, border=AMBER, alpha=12)
text(s, Inches(0.92), Inches(5.72), Inches(3.0), Inches(0.3),
     [("BILL OF MATERIALS", 11, AMBER, True)])
text(s, Inches(0.92), Inches(6.04), Inches(11.9), Inches(0.55),
     "%s   ·   %s   ·   %s   ·   %s"
     % (HARDWARE["pi"], HARDWARE["servos"], HARDWARE["power"], HARDWARE["cost"]),
     size=10.5, color=MUTED, spacing=1.2)

notes(s, """[SPEAKER 2 — 40 sec]
Hold the robot's hand up if you can.
"Each hand has five fingers and a wrist — six little servo motors — so twelve motors
in total, all driven by one board: the Adafruit PCA9685."
"Every finger is slightly different, so we measured and saved open-and-closed angles
for each one. That calibration is why the fist actually looks like a fist."
NOTE TO TEAM: fill the yellow bill-of-materials line (Pi model, servo model, power,
cost) in the CONFIG block before you print — those are not read off the code.
HAND OFF continues to the brain slide, still Speaker 2.""")


# ==========================================================================
# 06 - THE BRAIN  (controller + custom-LLM story)
# ==========================================================================
s = new_slide()
top = header(s, "04 · The brain", "Where Maxi actually thinks",
             "A cloud service on Railway — not a chip inside the robot", accent=PURPLE)

brain = [
    ("🚪", "Flask + Socket.IO", "The gateway the tablet talks to. Messages fly back and forth "
     "over a live connection.", CYAN),
    ("💬", "Groq language model", "llama-3.1-8b-instant turns the child's words into a friendly, "
     "kid-safe answer.", PURPLE),
    ("🧩", "Skills router", "Chat, maths, spelling, stories, quizzes, time & date — each is a "
     "small \"skill\".", GREEN),
    ("💾", "Long-term memory", "A tiny SQLite store so Maxi remembers a child's name and what "
     "they like, next time too.", PINK),
]
x = Inches(0.62)
for i, (emo, t, d, col) in enumerate(brain):
    px = Inches(0.62) + Emu(int(Inches(3.06) * (i % 2)))
    py = top + Emu(int(Inches(1.35) * (i // 2)))
    panel(s, px, py, Inches(2.9), Inches(1.25), alpha=14)
    badge(s, px + Inches(0.24), py + Inches(0.22), Inches(0.5), emo, col, glyph_size=15)
    text(s, px + Inches(0.86), py + Inches(0.28), Inches(1.9), Inches(0.32), t, size=12.5, bold=True, color=col)
    text(s, px + Inches(0.24), py + Inches(0.78), Inches(2.5), Inches(0.5), d, size=9.5, color=MUTED, spacing=1.12)

# custom-LLM honesty box (right)
panel(s, Inches(6.95), top, Inches(5.77), Inches(2.9), fill=PANEL_2, border=AMBER, alpha=12)
badge(s, Inches(7.25), top + Inches(0.28), Inches(0.6), "🎓", AMBER, glyph_size=18)
text(s, Inches(8.05), top + Inches(0.34), Inches(4.4), Inches(0.4),
     "We first trained our OWN model", size=16, bold=True, color=AMBER)
text(s, Inches(7.25), top + Inches(1.05), Inches(5.25), Inches(1.7),
     "At the start of the project we trained our own language model from scratch. But we "
     "could not run it — we did not have the computing power, and no funding to buy the "
     "hardware it needed. So we made an engineering choice: use a free cloud model instead, "
     "and put our effort into the robot. Same brain-in-the-cloud idea, zero licence cost.",
     size=12, color=MUTED, spacing=1.32)

# tech strip
panel(s, Inches(0.62), Inches(5.62), Inches(12.1), Inches(0.95), fill=PANEL, alpha=16)
tx = Inches(0.92)
for lbl, col, wd in [("Python", CYAN, 1.0), ("Flask 2.3", PURPLE, 1.05), ("Socket.IO", GREEN, 1.15),
                     ("groq 0.25", PINK, 1.1), ("Edge-TTS voice", AMBER, 1.5),
                     ("Railway cloud", CYAN, 1.35), ("free tier", GREEN, 0.95)]:
    pill(s, tx, Inches(5.94), lbl, col, w=Inches(wd), size=10.5)
    tx += Inches(wd) + Inches(0.14)

notes(s, """[SPEAKER 2 — 50 sec]
"Maxi's brain is not a chip inside the robot — it lives in the cloud."
Walk the four blocks fast: a gateway, a language model, a set of skills, and a memory.
THEN slow down for the yellow box — this is a highlight, tell it as a story:
"We actually trained our OWN model first. We were proud of it. But we could not run it —
we did not have the compute, and no funding to buy the hardware. So we made a real
engineering decision: use a free cloud model, and spend our time on the robot itself."
Judges and technical guests respect that honesty. HAND OFF stays with Speaker 2.""")


# ==========================================================================
# 07 - THE SENSES
# ==========================================================================
s = new_slide(orbs=((0.2, 0.18, CYAN), (0.82, 0.8, PINK)))
top = header(s, "05 · The senses", "How Maxi takes the world in", accent=CYAN)

senses = [
    ("🎤", "Microphone — hearing", "The tablet's mic is Maxi's main sense. The browser turns "
     "speech into text; that text is what the brain reads.", CYAN, "Web Speech API"),
    ("🔑", "Wake word — \"Hey Maxi\"", "A small on-device model (Vosk, ~41 MB) listens for its "
     "name, so it only wakes when spoken to.", PURPLE, "confidence ≥ 0.6"),
    ("👆", "Touch — the screen & head", "Tapping the screen, or Maxi's head, is a second input: "
     "a tap plays; a 3-second hold refreshes and resets it.", GREEN, "tablet touchscreen"),
]
x = Inches(0.62)
for emo, t, d, col, tag in senses:
    panel(s, x, top, Inches(3.86), Inches(3.5), alpha=14)
    strip = rect(s, x, top, Inches(3.86), Inches(0.075), col, MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    badge(s, x + Inches(0.36), top + Inches(0.42), Inches(0.74), emo, col, glyph_size=22)
    text(s, x + Inches(0.36), top + Inches(1.44), Inches(3.2), Inches(0.4), t, size=15.5, bold=True, color=col)
    text(s, x + Inches(0.36), top + Inches(1.96), Inches(3.25), Inches(1.15), d,
         size=11.5, color=MUTED, spacing=1.3)
    pill(s, x + Inches(0.36), top + Inches(3.0), tag, col, w=Inches(2.0), size=10)
    x += Inches(4.04)

panel(s, Inches(0.62), Inches(5.95), Inches(12.1), Inches(0.72), fill=PANEL_2, border=AMBER, alpha=12)
badge(s, Inches(0.92), Inches(6.06), Inches(0.5), "⚠️", AMBER, glyph_size=14)
text(s, Inches(1.6), Inches(6.08), Inches(10.9), Inches(0.5),
     [("Being honest:  ", 12.5, AMBER, True),
      ("Maxi has no camera and no distance sensors yet. Hearing is its main sense — giving "
       "it eyes is on our roadmap.", 12.5, INK, False)], spacing=1.2)

notes(s, """[SPEAKER 2 — 35 sec]
"Maxi mostly senses the world by listening."
Card 1: the tablet mic + speech-to-text is the main sense. Card 2: a small model on the
tablet listens just for the wake word, so it is not recording everything. Card 3: touch —
tapping the screen or its head.
Read the honesty strip out loud: "No camera yet — hearing is its main sense. Eyes are on
the roadmap." Do NOT pretend it has sensors it does not.
HAND OFF: "%s will show you how all of that turns into a decision." """
      % (PRESENTERS[2] if len(PRESENTERS) > 2 else PRESENTERS[0]))


# ==========================================================================
# 08 - HOW IT THINKS  (state machine)
# ==========================================================================
s = new_slide()
top = header(s, "06 · How it thinks", "A state machine, not a tangle of if-statements",
             "Maxi is always in exactly one state, and moves between them on clear events", accent=GREEN)

states = [("IDLE", "waiting", DIM), ("LISTENING", "hearing you", CYAN),
          ("THINKING", "asking the brain", PURPLE), ("SPEAKING", "talking + moving", GREEN)]
x = Inches(0.62)
for i, (nm, d, col) in enumerate(states):
    panel(s, x, top, Inches(2.72), Inches(1.5), alpha=16)
    text(s, x + Inches(0.26), top + Inches(0.28), Inches(2.3), Inches(0.4), nm, size=17, bold=True, color=col)
    text(s, x + Inches(0.26), top + Inches(0.82), Inches(2.3), Inches(0.5), d, size=12, color=MUTED)
    if i < 3:
        a = rect(s, x + Inches(2.78), top + Inches(0.55), Inches(0.24), Inches(0.36),
                 STROKE, MSO_SHAPE.CHEVRON)
        set_alpha(a, 45)
    x += Inches(3.06)
text(s, Inches(0.62), top + Inches(1.62), Inches(12.1), Inches(0.32),
     "…then back to IDLE. It is wake-gated: it only listens after you say \"Hey Maxi\" or tap — never on its own.",
     size=11.5, color=MUTED, italic=True)

rules = [
    ("🛡️", "It ignores its own voice", "Speech is accepted ONLY while the state is LISTENING. "
     "That one rule stops Maxi from hearing itself and answering itself in a loop.", CYAN),
    ("✋", "Answer 0–10 → real fingers", "If a maths answer is ten or less, Maxi counts it out on "
     "its hands. Bigger than ten, it shows it on the screen instead.", PURPLE),
    ("🎚️", "Smooth, not snappy", "Servos don't jump. Each move eases in and out with a cubic "
     "curve, ~67 updates a second — so movement looks natural, not robotic.", GREEN),
]
x = Inches(0.62)
for emo, t, d, col in rules:
    panel(s, x, Inches(4.35), Inches(3.86), Inches(2.15), alpha=14)
    badge(s, x + Inches(0.32), Inches(4.6), Inches(0.6), emo, col, glyph_size=17)
    text(s, x + Inches(1.06), Inches(4.66), Inches(2.7), Inches(0.35), t, size=13.5, bold=True, color=col)
    text(s, x + Inches(0.32), Inches(5.32), Inches(3.3), Inches(1.05), d,
         size=11.5, color=MUTED, spacing=1.28)
    x += Inches(4.04)

notes(s, """[SPEAKER 3 — 45 sec]
"People expect a robot to be a giant pile of if-statements. Ours is not. It is a state
machine — it is always in exactly ONE of four states."
Trace the top row: idle, listening, thinking, speaking, back to idle.
Then the three rules — pick the first one to dwell on: "Maxi only accepts what it hears
while it is in the listening state. That single rule is what stops it from hearing its own
voice and talking to itself forever — which is exactly the bug we hit first."
Mention: not PID, no motors racing — smooth cubic easing on the servos.""")


# ==========================================================================
# 09 - THE CODE  (real excerpts)
# ==========================================================================
s = new_slide(orbs=((0.2, 0.22, PURPLE), (0.8, 0.78, GREEN)))
top = header(s, "07 · The code", "Two real pieces of Maxi, straight from the repo", accent=PURPLE)

# excerpt 1 - state machine echo guard
panel(s, Inches(0.62), top, Inches(6.0), Inches(2.75), fill=PANEL, alpha=14)
text(s, Inches(0.92), top + Inches(0.2), Inches(5.4), Inches(0.35),
     "1 · The brain ignores its own echo", size=14, bold=True, color=CYAN)
code1 = panel(s, Inches(0.92), top + Inches(0.66), Inches(5.4), Inches(1.5), fill=BG_A, border=STROKE)
text(s, Inches(1.08), top + Inches(0.8), Inches(5.1), Inches(1.3),
     "# only accept speech while LISTENING\nif self.session.phase != Phase.LISTENING:\n    return   # <- drop Maxi's own voice",
     size=11, color=GREEN, font=MONO, spacing=1.25)
text(s, Inches(0.92), top + Inches(2.24), Inches(5.4), Inches(0.4),
     "Three lines that killed the \"robot talks to itself\" bug.", size=10.5, color=MUTED, italic=True)

# excerpt 2 - servo mapping
panel(s, Inches(6.82), top, Inches(5.9), Inches(2.75), fill=PANEL, alpha=14)
text(s, Inches(7.12), top + Inches(0.2), Inches(5.3), Inches(0.35),
     "2 · Angle → servo pulse → movement", size=14, bold=True, color=PINK)
code2 = panel(s, Inches(7.12), top + Inches(0.66), Inches(5.3), Inches(1.5), fill=BG_A, border=STROKE)
text(s, Inches(7.28), top + Inches(0.78), Inches(5.0), Inches(1.35),
     "pulse = 500 + (angle/180)*(2500-500)\nduty  = int(pulse / 20000 * 65535)\npca.channels[ch].duty_cycle = duty",
     size=10.5, color=AMBER, font=MONO, spacing=1.25)
text(s, Inches(7.12), top + Inches(2.24), Inches(5.3), Inches(0.4),
     "0–180° becomes a 500–2500 µs pulse at 50 Hz.", size=10.5, color=MUTED, italic=True)

# annotation strip
panel(s, Inches(0.62), Inches(5.5), Inches(12.1), Inches(1.05), fill=PANEL_2, border=PURPLE, alpha=12)
badge(s, Inches(0.95), Inches(5.76), Inches(0.56), "🔎", PURPLE, glyph_size=16)
text(s, Inches(1.66), Inches(5.72), Inches(10.8), Inches(0.75),
     [("No magic, no AI here — just clear rules.  ", 13.5, PURPLE, True),
      ("The left snippet is the safety rule for listening; the right one is the maths that "
       "turns an angle into an electrical pulse the servo understands. We wrote both.",
       13, INK, False)], spacing=1.25)

notes(s, """[SPEAKER 3 — 45 sec]
Do NOT read the code character by character. Explain each block in one breath.
LEFT: "A robot that talks can hear itself. These three lines say: only listen to new
speech while we are actually in the listening state. That fixed our first big bug."
RIGHT: "A servo does not understand 'ninety degrees'. It understands an electrical pulse.
This is the little bit of maths that converts an angle into a pulse width — 0 to 180
degrees becomes a pulse between 500 and 2500 microseconds."
Bottom line: this part is not AI, it is ordinary programming, and we wrote it ourselves.""")


# ==========================================================================
# 10 - WIRING  (built from shapes)
# ==========================================================================
s = new_slide()
top = header(s, "08 · Wiring", "How the three parts talk to each other", accent=CYAN)

nodes = [
    ("📱", "Tablet", "face · mic · speaker", CYAN),
    ("☁️", "Cloud brain", "Flask + Socket.IO\non Railway", PURPLE),
    ("🍓", "Raspberry Pi", "hand controller\nFlask REST · :5001", PINK),
    ("🎛️", "PCA9685", "16-ch PWM driver", GREEN),
    ("🤲", "12 servos", "fingers + wrists", AMBER),
]
link_labels = ["Socket.IO", "REST · X-API-Key", "I²C · SCL/SDA", "PWM · 50 Hz"]

node_w = Inches(2.16)
gap = Inches(0.29)
row_y = top + Inches(0.5)
nx = Inches(0.62)
for i, (emo, t, d, col) in enumerate(nodes):
    panel(s, nx, row_y, node_w, Inches(1.82), alpha=16)
    badge(s, nx + Emu(int((node_w - Inches(0.7)) / 2)), row_y + Inches(0.22), Inches(0.7), emo, col, glyph_size=20)
    text(s, nx, row_y + Inches(1.02), node_w, Inches(0.32), t, size=13.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    text(s, nx, row_y + Inches(1.36), node_w, Inches(0.44), d, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.05)
    if i < len(nodes) - 1:
        lx = nx + node_w
        a = rect(s, lx + Inches(0.02), row_y + Inches(0.7), gap - Inches(0.04), Inches(0.28),
                 STROKE, MSO_SHAPE.CHEVRON)
        set_alpha(a, 55)
        gap_center = lx + Emu(int(gap / 2))
        text(s, gap_center - Inches(0.9), top + Inches(0.04), Inches(1.8), Inches(0.4),
             link_labels[i], size=9, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.0)
    nx += node_w + gap

# legend rows (kept above the footer)
legend = [("Tablet ⇄ Brain", "Socket.IO messages over HTTP polling (Railway's server can't do raw WebSockets yet).", CYAN),
          ("Brain ⇄ Pi", "The brain POSTs commands like /show_number to the Pi, signed with a secret X-API-Key.", PINK),
          ("Pi ⇄ Servos", "The Pi speaks I²C to the PCA9685, which sends a 50 Hz PWM pulse to each of the 12 servos.", GREEN)]
y = row_y + Inches(2.1)
for t, d, col in legend:
    panel(s, Inches(0.62), y, Inches(12.1), Inches(0.72), fill=PANEL, alpha=14)
    rect(s, Inches(0.62), y, Inches(0.06), Inches(0.72), col)
    text(s, Inches(0.92), y + Inches(0.2), Inches(2.6), Inches(0.4), t, size=12.5, bold=True, color=col)
    text(s, Inches(3.4), y + Inches(0.11), Inches(9.1), Inches(0.55), d, size=11, color=MUTED, spacing=1.15,
         anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.8)

notes(s, """[SPEAKER 3 — 35 sec]
Trace the chain left to right with your finger:
"Tablet talks to the cloud brain. The brain sends hand commands to the Raspberry Pi.
The Pi speaks I²C to the servo board, and the board drives the twelve motors."
Name one real detail so it lands as real engineering: the Pi commands are signed with a
secret key, and the servos run on a 50 Hz pulse.
HAND OFF: "%s will show you everything Maxi can actually do — then the live demo." """
      % (PRESENTERS[3] if len(PRESENTERS) > 3 else PRESENTERS[0]))


# ==========================================================================
# 11 - WHAT IT CAN DO  (photo robot_3)
# ==========================================================================
s = new_slide(orbs=((0.18, 0.2, GREEN), (0.84, 0.8, CYAN)))
top = header(s, "09 · What it can do", "Maxi's skills today", accent=GREEN)

skills = [
    ("💬", "Talk & tutor", PURPLE), ("🔢", "Maths on fingers", CYAN),
    ("🔤", "Spelling", GREEN), ("📖", "Storytelling", PINK),
    ("❓", "Quizzes", AMBER), ("🕐", "Time & date", CYAN),
    ("💾", "Remembers you", PURPLE), ("😊", "Face feelings", GREEN),
    ("✌️", "Hand gestures", PINK), ("🗣️", "\"Hey Maxi\" wake", AMBER),
]
for i, (emo, t, col) in enumerate(skills):
    col_i = i % 2
    row_i = i // 2
    px = Inches(0.62) + Emu(int(Inches(2.78) * col_i))
    py = top + Emu(int(Inches(0.92) * row_i))
    panel(s, px, py, Inches(2.6), Inches(0.78), fill=PANEL, alpha=16)
    badge(s, px + Inches(0.2), py + Inches(0.17), Inches(0.46), emo, col, glyph_size=14)
    text(s, px + Inches(0.78), py + Inches(0.22), Inches(1.75), Inches(0.4), t, size=12, bold=True, color=col)

# photo on the right (clear of the tile columns)
framed_picture(s, PHOTO_3, Inches(6.35), top, Inches(6.37), Inches(3.75),
               border=GREEN, placeholder="[ robot_3.png — Maxi in action ]")

text(s, Inches(0.62), Inches(6.5), Inches(12.1), Inches(0.35),
     "Gestures the hands know: fist · peace · wave · count 1–5. Maths ≤ 10 is counted out on real fingers.",
     size=11.5, color=MUTED, italic=True)

notes(s, """[SPEAKER 4 — 35 sec]
Do not read all ten tiles. Pick three or four and say them like a child would ask:
"Hey Maxi, what's seven take away two? — and it counts five on its fingers." "Tell me a
story." "How do you spell banana?" "What time is it?"
Mention the one that always gets a reaction: "It remembers your name between sessions."
Then: "Rather than list them — let me show you." → advance to the demo.""")


# ==========================================================================
# 12 - TESTING
# ==========================================================================
s = new_slide()
top = header(s, "10 · Testing", "How we know it actually works", accent=PURPLE)

cards = [
    ("✅", "Automated tests", "We wrote tests that run with no robot attached — the voice engine, "
     "the memory, the maths and the interrupt logic all have their own checks that must pass.", CYAN),
    ("🎭", "Simulation mode", "The whole brain runs on a plain laptop with the hands simulated, so "
     "we could build and debug without the physical robot in front of us.", PURPLE),
    ("🎚️", "Calibration screen", "A settings page lets us drag each joint to its real open and "
     "closed angle and save it — that is where the per-finger numbers come from.", GREEN),
    ("📱", "Real-device passes", "Then we test on the actual tablet and hands, and tune the "
     "thresholds — wake sensitivity, timing — until it feels right.", PINK),
]
x, y = Inches(0.62), top
for i, (emo, t, d, col) in enumerate(cards):
    if i == 2:
        x, y = Inches(0.62), top + Inches(1.8)
    panel(s, x, y, Inches(6.0), Inches(1.62), alpha=14)
    badge(s, x + Inches(0.3), y + Inches(0.3), Inches(0.6), emo, col, glyph_size=17)
    text(s, x + Inches(1.06), y + Inches(0.24), Inches(4.7), Inches(0.35), t, size=14, bold=True, color=col)
    text(s, x + Inches(1.06), y + Inches(0.66), Inches(4.75), Inches(0.85), d, size=11, color=MUTED, spacing=1.25)
    x += Inches(6.1)

panel(s, Inches(0.62), Inches(5.72), Inches(12.1), Inches(0.9), fill=PANEL_2, border=CYAN, alpha=12)
text(s, Inches(1.0), Inches(5.92), Inches(11.4), Inches(0.55),
     [("Honest status:  ", 13, CYAN, True),
      ("the software is well tested; the physical hands still need the most tuning. We are "
       "up-front about what is solid and what is still rough.", 13, INK, False)], spacing=1.2)

notes(s, """[SPEAKER 4 — 25 sec]
"How do you know it works and we are not just hoping? Two ways."
Point left: "We wrote tests that run without the robot — the brain checks itself." Point
right-top: "The whole thing runs in a simulation on a laptop, so we did not need the robot
to build it." Calibration: "this screen is where the real finger angles come from."
Read the honest status line — software solid, hands still need tuning. Then set up the demo:
"Speaking of the hands — let's see them."  """)


# ==========================================================================
# 13 - LIVE DEMO
# ==========================================================================
s = new_slide(orbs=((0.25, 0.25, PINK), (0.75, 0.72, CYAN), (0.5, 0.9, PURPLE)))
_slide_no["n"] += 1

badge(s, Inches(6.16), Inches(1.15), Inches(1.0), "🚀", PINK, glyph_size=34)
text(s, Inches(0.6), Inches(2.4), Inches(12.13), Inches(1.0), "LIVE DEMO",
     size=60, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(3.42), Inches(12.13), Inches(0.5),
     "We say \"Hey Maxi\", ask a maths question, and its fingers answer",
     size=18, color=CYAN, align=PP_ALIGN.CENTER)

steps = [("1", "Say \"Hey Maxi\"", CYAN), ("2", "Ask a sum ≤ 10", PURPLE),
         ("3", "Watch the fingers", GREEN), ("4", "Try a word / story", PINK)]
x = Inches(1.35)
for num, lbl, col in steps:
    panel(s, x, Inches(4.35), Inches(2.62), Inches(1.0), fill=PANEL_2, alpha=20)
    text(s, x + Inches(0.24), Inches(4.5), Inches(0.55), Inches(0.4), num, size=22, bold=True, color=col)
    text(s, x + Inches(0.82), Inches(4.6), Inches(1.7), Inches(0.6), lbl, size=11.5, color=INK, spacing=1.1)
    x += Inches(2.77)

text(s, Inches(0.6), Inches(5.62), Inches(12.13), Inches(0.35),
     "If anything misbehaves: we do NOT debug on stage — we cut to the photo and keep talking.",
     size=12.5, color=AMBER, align=PP_ALIGN.CENTER, italic=True)

footer_only(s)

notes(s, """[SPEAKER 4 — 100 sec. THE MOST IMPORTANT SLIDE.]

PRE-FLIGHT CHECKLIST (do ALL of this BEFORE walking on stage):
  - Batteries: tablet fully charged; servo power pack fully charged; SPARE batteries in pocket.
  - Power the robot ON and let the hands finish their startup twitch BEFORE the talk.
  - Pair / connect the tablet to the brain and confirm the status says Connected.
  - Do ONE full dry run backstage: "Hey Maxi, what is three plus two?" and watch the fingers.
  - Open the chat screen already; wake word ON (or plan to tap the mic).
  - Clear the table around the hands — nothing for the fingers to catch on.
  - Volume up. Wi-Fi / hotspot confirmed (the brain is in the cloud — no internet, no Maxi).

ON STAGE (calm, slow):
  1. "Hey Maxi." Wait for the face to react and the beep. Do not rush it.
  2. "What is three plus two?" Let it think, speak, and COUNT FIVE ON ITS FINGERS. Stay quiet
     and let the audience watch the hand move — that silence is the whole show.
  3. Then a word: "How do you spell banana?" — it spells it out, letter by letter.
  4. Optional if time: "Tell me a story." Cut it off after one line.

IF IT MISBEHAVES — the recovery, rehearsed:
  Do NOT debug on stage. Say, warmly: "Maxi is being a little shy — here it is working
  earlier," then Alt-Tab / advance to the robot photo and KEEP TALKING. Tap-and-hold Maxi's
  head for 3 seconds to reset it while you speak; if it comes back, try once more. If not,
  move on. Nobody remembers a hiccup; they remember panic. So do not panic.""")


# ==========================================================================
# 14 - CHALLENGES & LESSONS
# ==========================================================================
s = new_slide()
top = header(s, "11 · What we learned", "The parts that did not go smoothly", accent=AMBER)

items = [
    ("🎓", "Our own model wouldn't run", "We trained one first, but had no compute and no funding "
     "for the hardware — so we moved to a free cloud model.", AMBER),
    ("🔁", "It talked to itself", "The robot heard its own voice and answered itself in a loop. "
     "We fixed it by only listening in the LISTENING state.", RED),
    ("🔊", "The mic beeped non-stop", "Always-on listening beeped on every phrase. We switched to "
     "wake-word + tap-to-talk so it is quiet until spoken to.", CYAN),
    ("🔑", "A tool went paid overnight", "Our first wake-word engine started requiring a company "
     "account, so we swapped in a free, offline one (Vosk).", PURPLE),
    ("💾", "The cloud forgot the child", "Railway wipes its disk on every deploy, so Maxi's memory "
     "reset. We attached a permanent volume to keep it.", GREEN),
    ("🔤", "Spelling came out garbled", "The voice ran letters together. With no SSML available we "
     "found a trick — \"B! A! N!\" — that forces clean letters.", PINK),
]
x, y = Inches(0.62), top
for i, (emo, t, d, col) in enumerate(items):
    if i == 3:
        x, y = Inches(0.62), top + Inches(2.28)
    panel(s, x, y, Inches(3.92), Inches(2.1), alpha=14)
    badge(s, x + Inches(0.26), y + Inches(0.26), Inches(0.56), emo, col, glyph_size=16)
    text(s, x + Inches(0.96), y + Inches(0.32), Inches(2.85), Inches(0.6), t, size=13, bold=True, color=col)
    text(s, x + Inches(0.26), y + Inches(1.0), Inches(3.42), Inches(1.0), d,
         size=11, color=MUTED, spacing=1.25)
    x += Inches(4.1)

notes(s, """[SPEAKER 1 — 30 sec]
You are back. Do NOT read six boxes. Tell TWO as short, real stories:
1. "We trained our own model and could not run it — no compute, no funding. Choosing to
   let that go and use the cloud was the hardest, most grown-up decision we made."
2. "For a whole afternoon, the robot kept interrupting itself — it could hear its own
   voice. The fix was one rule about when it is allowed to listen."
End with the line that ties it together: "The parts that broke taught us more than the
parts that worked." """)


# ==========================================================================
# 15 - LIMITS & WHAT'S NEXT
# ==========================================================================
s = new_slide(orbs=((0.18, 0.2, RED), (0.85, 0.78, PURPLE)))
top = header(s, "12 · Limits & next", "Being honest about what Maxi is — and isn't", accent=RED)

panel(s, Inches(0.62), top, Inches(6.0), Inches(3.7), fill=PANEL, border=RED, alpha=14)
badge(s, Inches(0.92), top + Inches(0.28), Inches(0.58), "⚠️", RED, glyph_size=17)
text(s, Inches(1.66), top + Inches(0.34), Inches(4.4), Inches(0.4), "What it can't do yet",
     size=16, bold=True, color=RED)
lim = [
    "It needs the internet — the brain is in the cloud.",
    "It talks to one child at a time, not a group.",
    "Hearing only — no camera, no eyes yet.",
    "The brain is a shared cloud model, not our own.",
    "The hands need calibrating; a servo can mis-move.",
    "It is a learning aid, not a certified teacher.",
]
y = top + Inches(0.98)
for l in lim:
    dot = rect(s, Inches(0.98), y + Inches(0.09), Inches(0.09), Inches(0.09), RED, MSO_SHAPE.OVAL)
    text(s, Inches(1.24), y, Inches(5.15), Inches(0.5), l, size=11.5, color=MUTED, spacing=1.2)
    y += Inches(0.44)

panel(s, Inches(6.92), top, Inches(5.8), Inches(3.7), fill=PANEL, border=GREEN, alpha=14)
badge(s, Inches(7.22), top + Inches(0.28), Inches(0.58), "🔭", GREEN, glyph_size=17)
text(s, Inches(7.96), top + Inches(0.34), Inches(4.4), Inches(0.4), "Where we'd take it next",
     size=16, bold=True, color=GREEN)
nxt = [
    ("Run our own model", "when we have the compute and the funding for it"),
    ("Give it eyes", "a camera, so it can see the child and the page"),
    ("More of the body", "arms and a moving head, not just hands"),
    ("Work offline", "so it runs without the internet"),
    ("Speak local languages", "Hausa and Kanuri, not only English"),
    ("Teach a whole group", "more than one child at once"),
]
y = top + Inches(0.98)
for t, d in nxt:
    dot = rect(s, Inches(7.28), y + Inches(0.09), Inches(0.09), Inches(0.09), GREEN, MSO_SHAPE.OVAL)
    text(s, Inches(7.54), y, Inches(4.95), Inches(0.5),
         [("%s — " % t, 11.5, INK, True), (d, 11.5, MUTED, False)], spacing=1.2)
    y += Inches(0.44)

panel(s, Inches(0.62), Inches(6.05), Inches(12.1), Inches(0.8), fill=PANEL_2, border=PURPLE, alpha=12)
text(s, Inches(1.0), Inches(6.24), Inches(11.4), Inches(0.5),
     [("The point was never a finished product. ", 14, PURPLE, True),
      ("It was to build something real that could one day help a child learn — and to prove to "
       "ourselves that we could.", 14, INK, False)], spacing=1.2)

notes(s, """[SPEAKER 2 — 30 sec]
Say the left box and mean it: "This is a school project. It needs the internet, it works
with one child, it has no eyes yet, and the brain is a shared cloud model, not our own. It
is a learning aid, not a teacher." Owning that earns you respect.
Then three from the right box, fast — running our own model, giving it eyes, local
languages. Land on the purple line, then advance to thank you.""")


# ==========================================================================
# 16 - THANK YOU  /  Q&A
# ==========================================================================
s = new_slide(orbs=((0.22, 0.24, PURPLE), (0.8, 0.76, CYAN), (0.5, 0.92, PINK)))
_slide_no["n"] += 1

badge(s, Inches(6.16), Inches(1.15), Inches(1.0), "🤖", PURPLE, glyph_size=34)
text(s, Inches(0.6), Inches(2.35), Inches(12.13), Inches(1.0), "Thank you",
     size=58, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(3.35), Inches(12.13), Inches(0.5),
     "Six months ago most of us had never written a line of code.",
     size=17, color=CYAN, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(3.76), Inches(12.13), Inches(0.5),
     "Today we built a robot that listens, talks, and teaches on its fingers.",
     size=17, color=MUTED, align=PP_ALIGN.CENTER)

pill(s, Inches(5.06), Inches(4.6), "QUESTIONS?", CYAN, w=Inches(3.2), size=15)

text(s, Inches(0.6), Inches(5.35), Inches(12.13), Inches(0.7),
     "%s\n%s  ·  %s  ·  %s"
     % (", ".join(PRESENTERS), BOOTCAMP, COHORT, EVENT_DATE),
     size=12, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.35)

footer_only(s)

notes(s, """[ALL FOUR SPEAKERS — 20 sec, then Q&A]
Everyone comes to the front. One person says the two lines, then: "We are happy to take
questions." Decide NOW who answers what.

LIKELY QUESTIONS & PREPARED ANSWERS:
Q: "Is Maxi's brain running inside the robot?"
   -> No. The thinking happens on a cloud server; the tablet and Pi talk to it over the
      internet. That is also why it needs a connection.
Q: "You said you trained your own AI — why aren't you using it?"
   -> We trained one, but we did not have the computing power to run it, and no funding to
      buy the hardware. We chose to use a free cloud model and spend our time on the robot.
Q: "How is this different from a talking toy or ChatGPT?"
   -> A toy plays fixed recordings. ChatGPT is just text. Maxi listens, decides, speaks AND
      moves real fingers to teach — and it remembers the child.
Q: "How does it count on its fingers?"
   -> Each finger is a servo motor on a driver board. If a maths answer is ten or less, the
      software opens exactly that many fingers.
Q: "What was the hardest part?"
   -> (your real answer) Letting go of our own model; and the day the robot kept
      interrupting itself because it could hear its own voice.""")


# --------------------------------------------------------------------------
OUT = "Maxi_Robot_Presentation.pptx"
prs.save(OUT)
print("Saved %s  (%d slides)" % (OUT, len(prs.slides._sldIdLst)))
